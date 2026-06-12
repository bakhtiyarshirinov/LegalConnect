from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colors ──────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x0A, 0x0A, 0x0A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x6B, 0x6B, 0x6B)
LGREY  = RGBColor(0xF5, 0xF5, 0xF5)
MGREY  = RGBColor(0xE8, 0xE8, 0xE8)
ACCENT = RGBColor(0x2D, 0x2D, 0x2D)

W  = Cm(33.87)   # slide width
H  = Cm(19.05)   # slide height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    from pptx.util import Emu
    sh = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh


def txbox(slide, text, x, y, w, h,
          size=Pt(18), bold=False, color=BLACK,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = "Calibri"
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb


def slide_title(slide, title, y=Cm(0.9)):
    # accent bar
    add_rect(slide, Cm(1.2), y, Cm(0.35), Cm(1.1), fill=BLACK)
    txbox(slide, title,
          Cm(1.8), y - Cm(0.05), Cm(30), Cm(1.3),
          size=Pt(28), bold=True, color=BLACK)


def bg_white(slide):
    add_rect(slide, 0, 0, W, H, fill=WHITE)


def bg_black(slide):
    add_rect(slide, 0, 0, W, H, fill=BLACK)


def feature_list(slide, items, x, y, size=Pt(15), gap=Cm(0.65), color=BLACK):
    for i, item in enumerate(items):
        txbox(slide, item, x, y + i * gap, Cm(14), Cm(0.7),
              size=size, color=color)


def placeholder_box(slide, x, y, w, h, label="[Screenshot]"):
    add_rect(slide, x, y, w, h, fill=LGREY, line=MGREY, line_w=Pt(1.5))
    txbox(slide, label, x, y + h/2 - Cm(0.4), w, Cm(0.8),
          size=Pt(14), color=GREY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, text, text_size=Pt(14), bold=False,
         bg=LGREY, border=MGREY, tcolor=BLACK):
    add_rect(slide, x, y, w, h, fill=bg, line=border, line_w=Pt(1.2))
    txbox(slide, text, x + Cm(0.3), y + Cm(0.2), w - Cm(0.6), h - Cm(0.4),
          size=text_size, bold=bold, color=tcolor, align=PP_ALIGN.CENTER, wrap=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_black(sl)

# decorative accent line
add_rect(sl, Cm(1.5), Cm(5.5), Cm(4), Cm(0.15), fill=WHITE)

txbox(sl, "LegalConnect",
      Cm(1.5), Cm(6.0), Cm(30), Cm(3.2),
      size=Pt(72), bold=True, color=WHITE)

txbox(sl, "Azərbaycanda Hüquqi Xidmətlər Platforması",
      Cm(1.5), Cm(9.4), Cm(28), Cm(1.2),
      size=Pt(22), color=RGBColor(0xCC, 0xCC, 0xCC))

txbox(sl, "Diplom İşi  •  2026",
      Cm(1.5), Cm(16.2), Cm(20), Cm(0.8),
      size=Pt(14), color=GREY)

txbox(sl, "Bakhtiyar Shirinov",
      Cm(1.5), Cm(17.2), Cm(20), Cm(0.8),
      size=Pt(13), color=GREY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Problem")

pain_points = [
    "❌  Azərbaycanda vəkil tapmaq çətindir",
    "❌  Şəffaf qiymət siyasəti yoxdur",
    "❌  Görüş planlaması vaxt aparır",
]
for i, pp in enumerate(pain_points):
    y = Cm(3.5) + i * Cm(2.2)
    add_rect(sl, Cm(1.2), y, Cm(17.5), Cm(1.8), fill=LGREY, line=MGREY)
    txbox(sl, pp, Cm(1.7), y + Cm(0.35), Cm(17), Cm(1.1),
          size=Pt(17), color=BLACK)

# Right — big stat
add_rect(sl, Cm(21), Cm(3.5), Cm(11), Cm(9), fill=BLACK)
txbox(sl, "70%",
      Cm(21), Cm(4.0), Cm(11), Cm(4.5),
      size=Pt(90), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "azərbaycanlıların\nhüquqi yardıma\nehtiyacı var",
      Cm(21), Cm(8.5), Cm(11), Cm(2.8),
      size=Pt(15), color=RGBColor(0xCC, 0xCC, 0xCC),
      align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Həll")

# Center logo text
txbox(sl, "LegalConnect",
      Cm(9.5), Cm(7.5), Cm(15), Cm(2.2),
      size=Pt(40), bold=True, color=BLACK, align=PP_ALIGN.CENTER)

boxes = [
    (Cm(1.0),  Cm(3.5),  "🔍  Vəkil Axtarışı"),
    (Cm(18.8), Cm(3.5),  "📅  Onlayn Görüş"),
    (Cm(1.0),  Cm(12.0), "💬  Real-time Söhbət"),
    (Cm(18.8), Cm(12.0), "📹  Video Zəng"),
]
for bx, by, btxt in boxes:
    card(sl, bx, by, Cm(12.5), Cm(2.2), btxt,
         text_size=Pt(17), bold=True, bg=LGREY, border=MGREY)

# connecting lines (decorative dashes)
add_rect(sl, Cm(13.5), Cm(8.5), Cm(6.5), Cm(0.06), fill=MGREY)
add_rect(sl, Cm(13.5), Cm(10.0), Cm(6.5), Cm(0.06), fill=MGREY)
add_rect(sl, Cm(16.5), Cm(5.5), Cm(0.06), Cm(7.0), fill=MGREY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Texnologiyalar")

# Left column header
add_rect(sl, Cm(1.2), Cm(3.0), Cm(14.5), Cm(0.9), fill=BLACK)
txbox(sl, "Backend", Cm(1.2), Cm(3.0), Cm(14.5), Cm(0.9),
      size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

backend = [
    ".NET 9  •  Clean Architecture",
    "CQRS + MediatR",
    "PostgreSQL + EF Core 9",
    "SignalR  •  JWT Auth",
    "MailKit (Gmail SMTP)",
]
for i, item in enumerate(backend):
    bg = WHITE if i % 2 == 0 else LGREY
    add_rect(sl, Cm(1.2), Cm(3.9) + i * Cm(1.1), Cm(14.5), Cm(1.0), fill=bg, line=MGREY)
    txbox(sl, item, Cm(1.7), Cm(3.95) + i * Cm(1.1), Cm(14), Cm(0.85),
          size=Pt(15), color=BLACK)

# Right column header
add_rect(sl, Cm(17.5), Cm(3.0), Cm(15), Cm(0.9), fill=BLACK)
txbox(sl, "Frontend / Mobile", Cm(17.5), Cm(3.0), Cm(15), Cm(0.9),
      size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

frontend = [
    "React 19 + TypeScript + Vite",
    "React Native + Expo SDK 56",
    "Zustand + React Query",
    "Framer Motion",
    "Daily.co API (Video)",
]
for i, item in enumerate(frontend):
    bg = WHITE if i % 2 == 0 else LGREY
    add_rect(sl, Cm(17.5), Cm(3.9) + i * Cm(1.1), Cm(15), Cm(1.0), fill=bg, line=MGREY)
    txbox(sl, item, Cm(18.0), Cm(3.95) + i * Cm(1.1), Cm(14.5), Cm(0.85),
          size=Pt(15), color=BLACK)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Arxitektura")

layers = [
    ("API Layer",            "Controllers • Middleware • Program.cs"),
    ("Application Layer",   "CQRS Handlers • DTOs • Validators"),
    ("Domain Layer",        "Entities • Interfaces • Business Rules"),
    ("Infrastructure Layer","EF Core • Repositories • SignalR Hub"),
]
colors = [BLACK, ACCENT, RGBColor(0x44, 0x44, 0x44), RGBColor(0x66, 0x66, 0x66)]

for i, ((lname, ldesc), col) in enumerate(zip(layers, colors)):
    y = Cm(3.2) + i * Cm(2.6)
    add_rect(sl, Cm(3.5), y, Cm(26.5), Cm(2.0), fill=col)
    txbox(sl, lname, Cm(3.8), y + Cm(0.1), Cm(12), Cm(0.9),
          size=Pt(17), bold=True, color=WHITE)
    txbox(sl, ldesc, Cm(3.8), y + Cm(0.9), Cm(26), Cm(0.85),
          size=Pt(13), color=RGBColor(0xCC, 0xCC, 0xCC))
    # arrow between layers
    if i < 3:
        txbox(sl, "▼", Cm(16.2), y + Cm(2.0), Cm(2), Cm(0.6),
              size=Pt(18), color=GREY, align=PP_ALIGN.CENTER)

txbox(sl, "Clean Architecture  •  SOLID Prinsipləri  •  Dependency Inversion",
      Cm(3.5), Cm(15.5), Cm(27), Cm(0.8),
      size=Pt(13), color=GREY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Features Overview
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Funksionallıq")

col_data = [
    ("Müştəri", [
        "✓ Qeydiyyat + OTP",
        "✓ Vəkil axtarışı",
        "✓ Görüş planla",
        "✓ Real-time söhbət",
        "✓ Video görüş",
        "✓ Rəy sistemi",
        "✓ Bildirişlər",
    ]),
    ("Vəkil", [
        "✓ Profil idarəetməsi",
        "✓ Slot yaratma",
        "✓ Görüş təsdiqi",
        "✓ Statistikalar",
        "✓ Real-time söhbət",
        "✓ Video görüş",
    ]),
    ("Admin", [
        "✓ Vəkil təsdiqi",
        "✓ İstifadəçilər",
        "✓ Statistika",
    ]),
]

col_xs = [Cm(1.2), Cm(12.3), Cm(23.4)]
col_w = Cm(10.2)

for (hdr, rows), cx in zip(col_data, col_xs):
    add_rect(sl, cx, Cm(2.9), col_w, Cm(0.9), fill=BLACK)
    txbox(sl, hdr, cx, Cm(2.9), col_w, Cm(0.9),
          size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        bg = WHITE if i % 2 == 0 else LGREY
        add_rect(sl, cx, Cm(3.8) + i * Cm(1.0), col_w, Cm(0.95), fill=bg, line=MGREY)
        txbox(sl, row, cx + Cm(0.4), Cm(3.85) + i * Cm(1.0), col_w - Cm(0.5), Cm(0.8),
              size=Pt(13), color=BLACK)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Client Portal
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Müştəri Portalı")

features = [
    "✓  Vəkil axtarışı və filtrlər",
    "✓  Slot əsaslı görüş planlaması",
    "✓  Real-time söhbət",
    "✓  Bildirişlər sistemi",
    "✓  Rəy sistemi",
    "✓  Profil idarəetməsi",
    "✓  Çoxdilli (AZ/EN/RU)",
]
feature_list(sl, features, Cm(1.2), Cm(3.2), size=Pt(16), gap=Cm(1.8))

placeholder_box(sl, Cm(17.5), Cm(2.8), Cm(14.5), Cm(13.5), "[Screenshot]")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Lawyer Portal
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Vəkil Portalı")

features = [
    "✓  Profil idarəetməsi",
    "✓  Slot yaratma (cədvəl)",
    "✓  Görüş təsdiqi / ləğvi",
    "✓  Müştəri söhbəti",
    "✓  Statistikalar paneli",
    "✓  Avatar yükləmə",
    "✓  Çoxdilli (AZ/EN/RU)",
]
feature_list(sl, features, Cm(1.2), Cm(3.2), size=Pt(16), gap=Cm(1.8))

placeholder_box(sl, Cm(17.5), Cm(2.8), Cm(14.5), Cm(13.5), "[Screenshot]")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Admin Panel
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Admin Paneli")

features = [
    "✓  Vəkil profilinin təsdiqi",
    "✓  İstifadəçi idarəetməsi",
    "✓  Platforma statistikası",
    "✓  Gözlənilən vəkillər siyahısı",
]
feature_list(sl, features, Cm(1.2), Cm(3.2), size=Pt(16), gap=Cm(1.8))

placeholder_box(sl, Cm(17.5), Cm(2.8), Cm(14.5), Cm(13.5), "[Screenshot]")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Mobile App
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Mobil Tətbiq")

txbox(sl, "React Native + Expo SDK 56",
      Cm(1.2), Cm(2.5), Cm(20), Cm(0.7),
      size=Pt(15), color=GREY, italic=True)

mob_cols = [
    ("📱\niOS / Android", "Expo Go ilə\ntest"),
    ("👤\nRole-based UI", "Müştəri + Vəkil\nayrı interfeys"),
    ("🌐\nÇoxdilli", "AZ / EN / RU\ndəstəyi"),
]
cx_start = Cm(1.2)
for i, (title, desc) in enumerate(mob_cols):
    cx = cx_start + i * Cm(10.8)
    add_rect(sl, cx, Cm(3.5), Cm(10.0), Cm(5.5), fill=LGREY, line=MGREY)
    txbox(sl, title, cx, Cm(3.7), Cm(10.0), Cm(2.2),
          size=Pt(18), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txbox(sl, desc, cx, Cm(5.9), Cm(10.0), Cm(1.5),
          size=Pt(13), color=GREY, align=PP_ALIGN.CENTER)

bottom_features = [
    "✓  Müştəri və Vəkil interfeysi",
    "✓  Görüş rezerv etmək",
    "✓  Vəkil axtarışı",
    "✓  Slot seçmə",
    "✓  Çat sistemi",
    "✓  Azərbaycan dili",
]
for i, f in enumerate(bottom_features):
    col = i // 3
    row = i % 3
    txbox(sl, f,
          Cm(1.2) + col * Cm(16),
          Cm(9.8) + row * Cm(1.5),
          Cm(15), Cm(1.2),
          size=Pt(14), color=BLACK)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Real-time Chat
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Real-time Söhbət")

txbox(sl, "Microsoft SignalR + WebSocket",
      Cm(1.2), Cm(2.5), Cm(25), Cm(0.7),
      size=Pt(15), color=GREY, italic=True)

# Flow diagram
flow_items = ["Müştəri", "SignalR Hub", "Vəkil"]
flow_colors = [LGREY, BLACK, LGREY]
flow_tcol   = [BLACK, WHITE, BLACK]
fx = Cm(2.0)
for i, (label, fc, tc) in enumerate(zip(flow_items, flow_colors, flow_tcol)):
    add_rect(sl, fx + i * Cm(8.5), Cm(3.5), Cm(7.5), Cm(1.8), fill=fc, line=MGREY)
    txbox(sl, label, fx + i * Cm(8.5), Cm(3.8), Cm(7.5), Cm(1.2),
          size=Pt(16), bold=True, color=tc, align=PP_ALIGN.CENTER)
    if i < 2:
        txbox(sl, "→", fx + Cm(7.5) + i * Cm(8.5), Cm(3.9), Cm(1.2), Cm(1.0),
              size=Pt(22), color=BLACK, align=PP_ALIGN.CENTER)

features = [
    "✓  Anlıq mesajlaşma",
    "✓  Fayl və şəkil göndərmə",
    "✓  Oxunmamış mesaj sayacı",
    "✓  Görüş təsdiqindən sonra aktivləşir",
    "✓  Çat axtarışı",
    "✓  Onlayn / oflayn status",
]
for i, f in enumerate(features):
    col = i // 3
    row = i % 3
    txbox(sl, f,
          Cm(1.2) + col * Cm(16),
          Cm(6.5) + row * Cm(1.6),
          Cm(15), Cm(1.2),
          size=Pt(15), color=BLACK)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — Video Meeting
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Video Görüş")

txbox(sl, "Daily.co API",
      Cm(1.2), Cm(2.5), Cm(20), Cm(0.7),
      size=Pt(15), color=GREY, italic=True)

# Flow
flow_steps = [
    "Görüş\ntəsdiqlənir",
    "Otaq\nyaradılır",
    "Email\ngöndərilir",
    "Görüşə\nqoşulur",
]
for i, step in enumerate(flow_steps):
    cx = Cm(1.2) + i * Cm(8.0)
    add_rect(sl, cx, Cm(3.5), Cm(7.2), Cm(2.2), fill=BLACK)
    txbox(sl, step, cx, Cm(3.6), Cm(7.2), Cm(2.0),
          size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        txbox(sl, "→", cx + Cm(7.2), Cm(3.9), Cm(0.9), Cm(1.0),
              size=Pt(20), color=BLACK, align=PP_ALIGN.CENTER)

features = [
    "✓  Avtomatik otaq yaradılması",
    "✓  Email bildirişi",
    "✓  Müştəri və vəkilə link",
    "✓  Görüş öncəsi 24 saat aktivləşir",
    "✓  Mobil dəstəyi",
]
for i, f in enumerate(features):
    txbox(sl, f, Cm(1.2), Cm(6.8) + i * Cm(1.6), Cm(28), Cm(1.2),
          size=Pt(15), color=BLACK)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — Security
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Təhlükəsizlik")

sec_items = [
    ("🔐", "JWT Token Auth",   "Bütün API endpoint\nqorunur"),
    ("📧", "Email OTP",        "Qeydiyyatda\ne-poçt doğrulaması"),
    ("🔒", "BCrypt Passwords", "Şifrələr hash\nkimi saxlanır"),
    ("👥", "Role-based Access","Admin / Vəkil /\nMüştəri icazələri"),
]

for i, (icon, title, desc) in enumerate(sec_items):
    cx = Cm(1.2) + i * Cm(8.0)
    add_rect(sl, cx, Cm(3.2), Cm(7.5), Cm(7.5), fill=LGREY, line=MGREY, line_w=Pt(1.5))
    txbox(sl, icon, cx, Cm(3.5), Cm(7.5), Cm(1.5),
          size=Pt(30), align=PP_ALIGN.CENTER)
    txbox(sl, title, cx, Cm(5.2), Cm(7.5), Cm(1.2),
          size=Pt(15), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txbox(sl, desc, cx, Cm(6.5), Cm(7.5), Cm(2.5),
          size=Pt(13), color=GREY, align=PP_ALIGN.CENTER)

txbox(sl, "Clean Architecture ilə təbəqəli müdafiə  •  HTTPS  •  Input Validation",
      Cm(1.2), Cm(15.8), Cm(31), Cm(0.8),
      size=Pt(13), color=GREY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Database
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Verilənlər Bazası")

txbox(sl, "PostgreSQL 16  +  Entity Framework Core 9",
      Cm(1.2), Cm(2.5), Cm(28), Cm(0.7),
      size=Pt(15), color=GREY, italic=True)

tables_left  = ["users", "lawyers", "appointments", "specializations", "availability_slots"]
tables_right = ["chats", "messages", "reviews", "notifications", "otp_codes"]

add_rect(sl, Cm(1.2), Cm(3.4), Cm(15), Cm(0.8), fill=BLACK)
txbox(sl, "Cədvəllər", Cm(1.2), Cm(3.4), Cm(15), Cm(0.8),
      size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, t in enumerate(tables_left):
    bg = WHITE if i % 2 == 0 else LGREY
    add_rect(sl, Cm(1.2), Cm(4.2) + i * Cm(1.05), Cm(15), Cm(1.0), fill=bg, line=MGREY)
    txbox(sl, t, Cm(1.6), Cm(4.25) + i * Cm(1.05), Cm(14.5), Cm(0.85),
          size=Pt(14), color=BLACK)

add_rect(sl, Cm(17.5), Cm(3.4), Cm(15), Cm(0.8), fill=BLACK)
txbox(sl, "Cədvəllər (davam)", Cm(17.5), Cm(3.4), Cm(15), Cm(0.8),
      size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, t in enumerate(tables_right):
    bg = WHITE if i % 2 == 0 else LGREY
    add_rect(sl, Cm(17.5), Cm(4.2) + i * Cm(1.05), Cm(15), Cm(1.0), fill=bg, line=MGREY)
    txbox(sl, t, Cm(17.9), Cm(4.25) + i * Cm(1.05), Cm(14.5), Cm(0.85),
          size=Pt(14), color=BLACK)

txbox(sl, "10 cədvəl  •  Docker Compose  •  EF Core Migrations",
      Cm(1.2), Cm(15.8), Cm(31), Cm(0.8),
      size=Pt(13), color=GREY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — Results
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_white(sl)
slide_title(sl, "Nəticələr")

stats = [
    ("4",    "Platforma\n(Web×3 + Mobile)"),
    ("10+",  "Əsas\nfunksiya"),
    ("3",    "Dil\ndəstəyi"),
    ("100%", "TypeScript\n(Frontend)"),
]
for i, (num, label) in enumerate(stats):
    cx = Cm(1.2) + i * Cm(8.0)
    add_rect(sl, cx, Cm(3.0), Cm(7.5), Cm(5.5), fill=BLACK)
    txbox(sl, num, cx, Cm(3.2), Cm(7.5), Cm(3.0),
          size=Pt(52), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, label, cx, Cm(6.2), Cm(7.5), Cm(2.0),
          size=Pt(13), color=RGBColor(0xCC, 0xCC, 0xCC),
          align=PP_ALIGN.CENTER)

bottom_stats = [
    ("Backend",  "50+ API endpoint"),
    ("Frontend", "100+ komponent"),
    ("Mobile",   "20+ ekran"),
]
for i, (lbl, val) in enumerate(bottom_stats):
    cx = Cm(3.0) + i * Cm(9.5)
    add_rect(sl, cx, Cm(10.0), Cm(8.5), Cm(2.8), fill=LGREY, line=MGREY, line_w=Pt(1.5))
    txbox(sl, lbl, cx, Cm(10.1), Cm(8.5), Cm(0.9),
          size=Pt(13), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txbox(sl, val, cx, Cm(11.0), Cm(8.5), Cm(1.2),
          size=Pt(18), bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 — Conclusion
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_black(sl)

add_rect(sl, Cm(1.5), Cm(5.5), Cm(4), Cm(0.15), fill=WHITE)

txbox(sl, "LegalConnect",
      Cm(1.5), Cm(6.0), Cm(30), Cm(3.0),
      size=Pt(72), bold=True, color=WHITE)

txbox(sl, "Azərbaycanda hüquqi xidmətləri rəqəmsallaşdırırıq",
      Cm(1.5), Cm(9.3), Cm(30), Cm(1.0),
      size=Pt(20), color=RGBColor(0xCC, 0xCC, 0xCC))

add_rect(sl, Cm(1.5), Cm(12.5), Cm(30), Cm(0.06), fill=RGBColor(0x33, 0x33, 0x33))

txbox(sl, "github.com/bakhtiyarshirinov/LegalConnect",
      Cm(1.5), Cm(13.2), Cm(30), Cm(0.8),
      size=Pt(14), color=GREY)

txbox(sl, "Təşəkkür edirik!",
      Cm(1.5), Cm(15.5), Cm(30), Cm(1.5),
      size=Pt(32), bold=True, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
out = "/Users/mac/LegalConnect/LegalConnect_Presentation.pptx"
prs.save(out)
print(f"Presentation created successfully! → {out}")
print(f"Total slides: {len(prs.slides)}")
